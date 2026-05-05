"""Build a standalone slide deck for Jiban + Poornika to review:
PIM Enablement architecture + Hybrid identity recommendation.

Output: partner-context/ZAFIN_PIM_ENABLEMENT_DECK_MAY2026.pptx

Run:  python scripts/build_pim_enablement_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- Theme ---------------------------------------------------------------
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
SUBTLE = RGBColor(0x60, 0x60, 0x60)
GREEN = RGBColor(0x10, 0x7C, 0x10)
AMBER = RGBColor(0xCA, 0x5A, 0x00)
RED = RGBColor(0xC4, 0x31, 0x4B)
LIGHT_BG = RGBColor(0xF3, 0xF7, 0xFB)
SOFT_GRAY = RGBColor(0xE5, 0xEA, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUT_PATH = (
    Path(__file__).resolve().parent.parent
    / "partner-context"
    / "ZAFIN_PIM_ENABLEMENT_DECK_MAY2026.pptx"
)


# ---- Helpers -------------------------------------------------------------
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, text, left, top, width, height, *, size=14, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, name="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = name
    return tb


def add_bullets(slide, items, left, top, width, height, *, size=12, color=DARK_TEXT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = f"•  {it}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb


def add_card(slide, x, y, w, h, *, fill=WHITE, border=AZURE_BLUE, border_pt=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = border; card.line.width = Pt(border_pt)
    return card


def add_pill(slide, x, y, w, h, text, *, fill, text_color=WHITE, size=11):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    pill.fill.solid(); pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = text_color; r.font.name = "Segoe UI"
    return pill


def add_header_bar(prs, slide, title, kicker=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = AZURE_BLUE
    add_text(slide, title, 0.5, 0.55, 12.5, 0.55, size=24, bold=True)
    if kicker:
        add_text(slide, kicker, 0.5, 1.05, 12.5, 0.35, size=12, color=SUBTLE)


def add_footer(slide, idx, total):
    add_text(slide,
             "Arturo Quiroga  ·  PSA  ·  Zafin engagement  ·  May 2026  ·  Internal review",
             0.5, 7.18, 9.5, 0.25, size=9, color=SUBTLE)
    add_text(slide, f"{idx} / {total}", 12.0, 7.18, 1.0, 0.25, size=9, color=SUBTLE, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def connector(slide, x1, y1, x2, y2, *, color=AZURE_BLUE, width_pt=1.5, arrow=True):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if arrow:
        # Try to add arrow head; python-pptx exposes via XML
        from pptx.oxml.ns import qn
        ln = line.line._get_or_add_ln()
        tail = ln.find(qn("a:tailEnd"))
        if tail is None:
            from lxml import etree
            tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med"); tail.set("len", "med")
    return line


# ---- Build ---------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 8

# === Slide 1: Title =======================================================
s = add_blank(prs)
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.4), prs.slide_height)
side.line.fill.background(); side.fill.solid(); side.fill.fore_color.rgb = AZURE_BLUE
add_text(s, "ZAFIN", 0.5, 1.3, 3.5, 0.8, size=44, bold=True, color=WHITE)
add_text(s, "PIM Enablement", 0.5, 2.1, 3.5, 0.6, size=22, color=WHITE)
add_text(s, "+ Hybrid identity model", 0.5, 2.55, 3.5, 0.5, size=16, color=WHITE)
add_text(s, "Architecture review", 4.9, 1.5, 8.0, 0.7, size=32, bold=True)
add_text(s, "Pre-read for the next internal alignment session", 4.9, 2.3, 8.0, 0.5, size=18, color=SUBTLE)
add_bullets(s, [
    "End-of-June 2026 milestone for Zafin",
    "Banking-compliance constraint: human-in-loop is mandatory",
    "Recommendation: ship PIM Enablement first; Self-Healing V1 second",
], 4.9, 3.5, 8.0, 2.0, size=14, color=DARK_TEXT)
add_text(s,
         "Arturo Quiroga  ·  PSA  ·  for Jiban Biswas + Poornika Dodamgodage  ·  May 4, 2026",
         4.9, 6.5, 8.0, 0.4, size=12, color=SUBTLE)
add_notes(s,
          "Pre-read deck for Jiban (new on the team) and Poornika before the next internal alignment "
          "session. Two topics: PIM Enablement architecture, and the hybrid identity recommendation "
          "that unblocks both PIM and Self-Healing V1.")

# === Slide 2: Context & priorities =======================================
s = add_blank(prs)
add_header_bar(prs, s, "Where we are", "Latest customer input: Apr 28 working session + May 1 use-case writeup")

# Two columns: numbers and decisions
add_text(s, "Numbers we now have", 0.5, 1.5, 6.0, 0.4, size=16, bold=True, color=AZURE_BLUE)
nums_card = add_card(s, 0.5, 1.95, 6.0, 2.6, fill=LIGHT_BG)
add_bullets(s, [
    "Operational volume: ~2,000–3,000 tickets / day",
    "Q2 target: 5–10% ticket automation",
    "Q3 target: 25% ticket automation",
    "Hard milestone: end of June 2026 for first MS-side delivery",
    "Constraint: banking compliance \u2192 human-in-loop is mandatory",
], 0.7, 2.1, 5.6, 2.3, size=13)

add_text(s, "Decisions taken", 6.8, 1.5, 6.0, 0.4, size=16, bold=True, color=AZURE_BLUE)
dec_card = add_card(s, 6.8, 1.95, 6.0, 2.6, fill=LIGHT_BG)
add_bullets(s, [
    "PIM Enablement = P1, end of June",
    "Infra & Ops Automation V1 = P2, end of June",
    "Customer Enablement = deferred",
    "Read-only RBAC pattern validated (Apr 14 + Apr 28)",
    "Ankit Desai off the team \u2192 Jiban Biswas joining",
], 7.0, 2.1, 5.6, 2.3, size=13)

# Bottom risk strip
risk = add_card(s, 0.5, 4.85, 12.3, 2.0, fill=WHITE, border=AMBER, border_pt=1.5)
add_text(s, "Open risks for the end-of-June milestone", 0.7, 4.95, 12.0, 0.4, size=14, bold=True, color=AMBER)
add_bullets(s, [
    "Two P1s on the same date \u2014 propose sequencing: PIM mid-June, Ops V1 demo end-of-June",
    "PIM trigger mechanism (Event Grid vs poll) and Teams output mechanism not yet picked",
    "Validation criteria for PIM owned by Zafin \u2014 not yet shared",
    "Optimus partner overlap with Self-Healing scope \u2014 boundary doc still outstanding",
    "No baseline metric / instrumentation plan for the Q2/Q3 ticket-automation %",
], 0.7, 5.35, 12.0, 1.5, size=12)
add_footer(s, 2, TOTAL)
add_notes(s,
          "Two-column reset of where we are. Numbers are the new news from Poornika's May 1 writeup. "
          "Decisions reflect Apr 28 + May 1. Risks are the things we should land in the next internal "
          "session before facing the customer again.")

# === Slide 3: PIM Enablement \u2014 the flow ====================================
s = add_blank(prs)
add_header_bar(prs, s, "PIM Enablement \u2014 end-to-end flow",
               "Agent validates and recommends. A human always approves in the PIM portal.")

# 6 step boxes across the slide
steps = [
    ("1\nUser raises\nPIM request", "Entra PIM"),
    ("2\nAgent reads\nrequest body", "MS Graph"),
    ("3\nValidate group\nmembership", "MS Graph"),
    ("4\nValidate Jira\nticket", "Jira MCP \u2713"),
    ("5\nApply Zafin\nrules", "Local logic"),
    ("6\nAdaptive Card\n\u2192 Teams", "Graph / webhook"),
]
y_steps = 1.7
box_w = 1.95
gap = 0.15
total_w = len(steps) * box_w + (len(steps) - 1) * gap
x_start = (13.333 - total_w) / 2
for i, (label, sub) in enumerate(steps):
    x = x_start + i * (box_w + gap)
    box = add_card(s, x, y_steps, box_w, 1.5, fill=LIGHT_BG, border=AZURE_BLUE, border_pt=1.5)
    add_text(s, label, x + 0.05, y_steps + 0.1, box_w - 0.1, 1.05, size=12, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, sub, x + 0.05, y_steps + 1.15, box_w - 0.1, 0.3, size=10, color=SUBTLE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        connector(s, x + box_w, y_steps + 0.75, x + box_w + gap, y_steps + 0.75)

# Output panel
out_y = 3.55
add_text(s, "Output \u2192 Adaptive Card posted to approver's Teams channel", 0.5, out_y, 12.5, 0.4, size=14, bold=True, color=AZURE_BLUE)
card_x = 0.5; card_w = 7.0; card_h = 3.2
add_card(s, card_x, out_y + 0.45, card_w, card_h, fill=WHITE, border=SUBTLE, border_pt=1.0)
add_text(s, "🔐  PIM Activation \u2014 Recommendation", card_x + 0.2, out_y + 0.55, card_w - 0.4, 0.35, size=13, bold=True)
add_bullets(s, [
    "Requester · role · scope · justification",
    "Linked Jira ticket: INC-12345",
    "Validation checklist (✓ / ✗ per rule)",
    "Verdict:  ELIGIBLE  ·  Confidence: High",
    "Recommended action: Approve",
    "Buttons: [Open in PIM Portal]  [View Ticket]  [View Evidence]",
], card_x + 0.2, out_y + 0.95, card_w - 0.4, 2.5, size=11)
# Note: NO Approve button on the card itself
add_pill(s, card_x + 0.2, out_y + card_h + 0.15, 4.0, 0.35,
         "No \u201CApprove\u201D button on the card", fill=AMBER, size=10)

# Right panel: why this is compliance-safe
rx = 8.0; rw = 5.0
add_card(s, rx, out_y + 0.45, rw, card_h, fill=LIGHT_BG, border=GREEN, border_pt=1.5)
add_text(s, "Why this is banking-compliance safe", rx + 0.2, out_y + 0.55, rw - 0.4, 0.35, size=13, bold=True, color=GREEN)
add_bullets(s, [
    "Agent never holds RoleAssignment.ReadWrite",
    "Agent never holds Approve permission on PIM",
    "Final action stays in PIM portal (already audited)",
    "Every recommendation logs reasoning + evidence to Jira (immutable trail)",
    "Read-only end-to-end \u2014 no new RBAC required for the agent itself",
], rx + 0.2, out_y + 0.95, rw - 0.4, 2.5, size=11)

add_footer(s, 3, TOTAL)
add_notes(s,
          "Six-step flow mapped to its integration points. The card design is the key compliance "
          "lever: there is no Approve button on the card, so the human must act in the PIM portal "
          "where it's already audited. Every check is itemized so the human approver gets the "
          "evidence they need without re-doing the work.")

# === Slide 4: Hybrid identity recommendation =============================
s = add_blank(prs)
add_header_bar(prs, s, "Identity model \u2014 recommendation",
               "Hybrid: MI for reads, OBO for writes  \u2014  defensible attribution for banking-regulated workloads")

card_w = 4.1; card_h = 3.6; gap = 0.2; y = 1.55
cards = [
    {
        "title": "Managed Identity only",
        "sub": "Agent acts with its own elevated rights",
        "pros": ["Simple to implement", "Works in async / scheduled flows", "Clean MI audit trail"],
        "cons": ["Confused-deputy risk", "Privilege held permanently", "Hard to attribute to a human"],
        "verdict": "Not acceptable for banking",
        "vc": RED, "border": SUBTLE,
    },
    {
        "title": "On-Behalf-Of only",
        "sub": "Agent acts under requesting user's token",
        "pros": ["Inherits user RBAC", "Action attributed to the human", "Clean compliance story"],
        "cons": ["Async / scheduled flows fail", "Token lifetime issues", "No proactive sweeps"],
        "verdict": "Partial \u2014 writes only",
        "vc": AMBER, "border": SUBTLE,
    },
    {
        "title": "Hybrid  (recommended)",
        "sub": "MI for reads, OBO for writes",
        "pros": ["Reads stay on validated MI pattern", "Every write carries human attribution", "Defensible audit trail per action"],
        "cons": ["Slightly more implementation complexity", "User identity must be in the request chain"],
        "verdict": "Recommended for banking compliance",
        "vc": GREEN, "border": GREEN,
    },
]
for i, c in enumerate(cards):
    x = 0.5 + i * (card_w + gap)
    add_card(s, x, y, card_w, card_h, border=c["border"], border_pt=2.0 if i == 2 else 1.0)
    title_color = c["border"] if i == 2 else DARK_TEXT
    add_text(s, c["title"], x + 0.2, y + 0.15, card_w - 0.4, 0.4, size=15, bold=True, color=title_color)
    add_text(s, c["sub"], x + 0.2, y + 0.55, card_w - 0.4, 0.4, size=11, color=SUBTLE)
    add_text(s, "Pros", x + 0.2, y + 1.0, card_w - 0.4, 0.3, size=11, bold=True, color=GREEN)
    add_bullets(s, c["pros"], x + 0.2, y + 1.3, card_w - 0.4, 1.0, size=10)
    add_text(s, "Cons", x + 0.2, y + 2.35, card_w - 0.4, 0.3, size=11, bold=True, color=RED)
    add_bullets(s, c["cons"], x + 0.2, y + 2.65, card_w - 0.4, 1.0, size=10)
for i, c in enumerate(cards):
    x = 0.5 + i * (card_w + gap)
    add_pill(s, x, y + card_h + 0.1, card_w, 0.4, c["verdict"], fill=c["vc"], size=12)

# Bottom mapping strip
yb = 5.85
add_card(s, 0.5, yb, 12.3, 1.15, fill=LIGHT_BG, border=AZURE_BLUE, border_pt=1.0)
add_text(s, "Applied to Zafin's two priority use cases", 0.7, yb + 0.07, 12.0, 0.35, size=13, bold=True, color=AZURE_BLUE)
mapping = [
    ("PIM Enablement (P1)", "MI only \u2014 read-only end-to-end",
     "Agent never writes; recommends only. Human approves in PIM portal."),
    ("Ops Automation V1 (P2)", "MI for reads + OBO for write actions",
     "restart / scale / rollback executed under requesting user's token. Dev subscription only."),
]
for i, (uc, model, why) in enumerate(mapping):
    x = 0.7 + i * 6.1
    add_text(s, uc, x, yb + 0.45, 6.0, 0.3, size=12, bold=True)
    add_text(s, model, x, yb + 0.7, 6.0, 0.25, size=11, bold=True, color=GREEN)
    add_text(s, why, x, yb + 0.92, 6.0, 0.3, size=10, color=SUBTLE)

add_footer(s, 4, TOTAL)
add_notes(s,
          "Lead with this. MI-only is off the table for banking. OBO-only breaks proactive flows. "
          "Hybrid is the only model that satisfies both. Applied: PIM Enablement is MI-only because "
          "it's read-only end-to-end. Ops Automation V1 uses MI for reads + OBO for writes.")

# === Slide 5: Integration points & gaps ==================================
s = add_blank(prs)
add_header_bar(prs, s, "Integration points  +  gaps to close",
               "What we have, what we need, and what's still unknown")

# Table-style rows
headers = ["#", "Integration", "Status", "Owner / next step"]
rows = [
    ("1", "Entra PIM trigger (webhook vs poll)", "GAP", "Spike week 1; fall back to poll", AMBER),
    ("2", "MS Graph PIM API \u2014 read request", "Standard \u2713", "Use Graph SDK", GREEN),
    ("3", "Group membership read (Graph)", "Standard \u2713", "getMemberGroups", GREEN),
    ("4", "Jira ticket read + audit comment", "Existing MCP \u2713", "From Jan PoC", GREEN),
    ("5", "Validation rules engine (Zafin criteria)", "Local logic", "Zafin owns the rule table", AMBER),
    ("6", "Teams output (Graph chat vs webhook)", "GAP", "Decide in next internal meeting", AMBER),
    ("7", "Audit trail (Jira comment + immutable log?)", "Partial", "Confirm format with Zafin compliance", AMBER),
    ("8", "PIM pending-request reads (Enterprise MCP gap)", "RESOLVED", "Hybrid: gap-filler `pim-mcp` live 2026-05-05 \u2014 see slide 6", GREEN),
]

y0 = 1.5
row_h = 0.55
header_h = 0.5
# Header row
header_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y0), Inches(12.3), Inches(header_h))
header_bg.fill.solid(); header_bg.fill.fore_color.rgb = AZURE_BLUE; header_bg.line.fill.background()
col_x = [0.6, 1.2, 7.0, 9.0]
col_w_pts = [0.6, 5.7, 1.9, 3.6]
for i, h in enumerate(headers):
    add_text(s, h, col_x[i], y0 + 0.1, col_w_pts[i], 0.35, size=12, bold=True, color=WHITE)

for i, (num, integ, status, next_step, color) in enumerate(rows):
    yr = y0 + header_h + i * row_h
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(yr), Inches(12.3), Inches(row_h))
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else SOFT_GRAY
    bg.line.color.rgb = SOFT_GRAY; bg.line.width = Pt(0.5)
    add_text(s, num, col_x[0], yr + 0.13, col_w_pts[0], 0.3, size=11, bold=True, color=SUBTLE)
    add_text(s, integ, col_x[1], yr + 0.13, col_w_pts[1], 0.3, size=11)
    # status pill
    add_pill(s, col_x[2], yr + 0.1, 1.7, 0.35, status, fill=color, size=10)
    add_text(s, next_step, col_x[3], yr + 0.13, col_w_pts[3], 0.3, size=11, color=SUBTLE)

add_footer(s, 5, TOTAL)
add_notes(s,
          "Three real gaps: PIM trigger mechanism, Teams output mechanism, and whether an "
          "MS-supported PIM MCP exists. Surface these in the meeting and assign owners. "
          "Don't let them slide \u2014 they are the schedule risk for end-of-June. "
          "UPDATE 2026-05-05: row #8 (PIM connector reality check) is RESOLVED via the hybrid "
          "approach \u2014 see next slide for details.")

# === Slide 6: Live validation update (May 5, 2026) =======================
s = add_blank(prs)
add_header_bar(prs, s, "Live validation update  \u00b7  2026-05-05",
               "Hybrid Enterprise MCP + gap-filler `pim-mcp` is working end-to-end in `MngEnvMCAP094150`")

# Left: what works
add_text(s, "What works today", 0.5, 1.5, 6.0, 0.4, size=16, bold=True, color=GREEN)
add_card(s, 0.5, 1.95, 6.0, 4.5, fill=WHITE, border=GREEN, border_pt=1.5)
add_bullets(s, [
    "`pim-mcp` Container App live (image 0.2.4)",
    "SSE \u2192 FastMCP \u2192 DefaultAzureCredential \u2192 IMDS \u2192 MI token",
    "Microsoft Graph LIST `roleAssignmentScheduleRequests` \u2192 HTTP 200",
    "`list_pending_pim_requests` returns normalized JSON for Foundry",
    "Three reproducible PowerShell scripts published:",
    "   \u2022  assign-pim-eligibility.ps1",
    "   \u2022  configure-pim-approval.ps1",
    "   \u2022  trigger-pim-activation.ps1",
    "End-to-end runbook: scripts/REPRODUCE.md",
], 0.7, 2.1, 5.6, 4.3, size=12)

# Right: upstream bugs found
add_text(s, "Upstream bugs filed (UPSTREAM_BUGS.md)", 6.8, 1.5, 6.0, 0.4, size=16, bold=True, color=AMBER)
add_card(s, 6.8, 1.95, 6.0, 4.5, fill=WHITE, border=AMBER, border_pt=1.5)
add_bullets(s, [
    "BUG-001 (High): Graph runtime requires `*.ReadWrite.Directory` for read-only LIST \u2014 docs claim Read.Directory works",
    "BUG-002 (Medium): Graph rejects $filter, $orderby AND $expand on this collection (server filters client-side)",
    "BUG-003 (Medium): Enterprise MCP exposes only `MCP.*Read.*` scopes \u2014 cannot satisfy BUG-001",
    "App-role propagation lag: 5\u201360+ min after grant (resource-side claim cache)",
], 7.0, 2.1, 5.6, 4.3, size=11)

# Bottom: compliance posture restatement
yb = 6.55
add_card(s, 0.5, yb, 12.3, 0.85, fill=LIGHT_BG, border=AZURE_BLUE, border_pt=1.0)
add_text(s, "Compliance posture unchanged", 0.7, yb + 0.07, 12.0, 0.3, size=12, bold=True, color=AZURE_BLUE)
add_text(s,
         "MI now holds RoleAssignmentSchedule.ReadWrite.Directory at runtime (BUG-001 workaround). "
         "Compensating control: `pim-mcp` server registers ZERO write tools \u2014 the elevated token is latent only. "
         "Documented in threat-model.md Residual Risk #4.",
         0.7, yb + 0.35, 12.0, 0.45, size=11, color=DARK_TEXT)

add_footer(s, 6, TOTAL)
add_notes(s,
          "Brand-new slide for the May 5 update. The gap-filler hybrid that was planned on slide 4 "
          "is now live and validated end-to-end against Microsoft Graph in our MngEnvMCAP094150 "
          "tenant. Three upstream Graph/Enterprise-MCP bugs were uncovered in the process \u2014 all "
          "documented and worked around. Compliance posture is unchanged because the server "
          "registers no write tool, even though the MI's token now carries ReadWrite scope at "
          "runtime (forced by BUG-001).")

# === Slide 7: Self-Healing V1 scope (lane A) =============================
s = add_blank(prs)
add_header_bar(prs, s, "Infrastructure & Operations Automation  \u2014  V1 scope",
               "Tightly bounded: lane A only, Dev subscription, user-approved")

# Left: in scope
add_text(s, "In scope for V1", 0.5, 1.5, 6.0, 0.4, size=16, bold=True, color=GREEN)
add_card(s, 0.5, 1.95, 6.0, 4.5, fill=WHITE, border=GREEN, border_pt=1.5)
add_bullets(s, [
    "Lane A only \u2014 ticket resolution automation",
    "2\u20133 high-volume Dev ticket categories (Zafin to nominate)",
    "Allowlisted write actions: restart  \u00b7  scale  \u00b7  rollback",
    "User approves each action before execution",
    "Hybrid identity: MI for reads, OBO for writes",
    "Dev subscription only \u2014 hard guardrail",
    "Before / after evidence captured to Jira (audit trail)",
    "Agent verifies outcome after action runs",
], 0.7, 2.1, 5.6, 4.3, size=13)

# Right: explicitly out
add_text(s, "Explicitly out of V1", 6.8, 1.5, 6.0, 0.4, size=16, bold=True, color=RED)
add_card(s, 6.8, 1.95, 6.0, 4.5, fill=WHITE, border=RED, border_pt=1.5)
add_bullets(s, [
    "Lane B \u2014 runbook execution engine (defer)",
    "Lane C \u2014 pattern-matched auto-remediation (V2)",
    "Config-map / secret edits",
    "RBAC or network-policy changes",
    "Anything cluster-scoped",
    "Cross-namespace operations",
    "Production subscriptions",
    "Proactive write actions (no user trigger)",
], 7.0, 2.1, 5.6, 4.3, size=13)

# Bottom rationale
add_text(s,
         "Rationale: end-of-June is 8 weeks away. Lane A with 2\u20133 actions is achievable, "
         "demoable, and de-risks the harder design questions before V2.",
         0.5, 6.55, 12.5, 0.45, size=12, color=SUBTLE, align=PP_ALIGN.CENTER)
add_footer(s, 7, TOTAL)
add_notes(s,
          "Frame V1 as lane A only. Explicitly list what is OUT to anchor expectations. "
          "Ask Zafin which 2\u20133 ticket categories represent the most volume; let the data drive "
          "the allowlist instead of guessing.")

# === Slide 8: Asks & next steps ==========================================
s = add_blank(prs)
add_header_bar(prs, s, "Asks  \u00b7  next steps", "What I need from this internal session")

# Left: asks
add_text(s, "Decisions I'm asking for", 0.5, 1.5, 6.0, 0.4, size=16, bold=True, color=AZURE_BLUE)
add_card(s, 0.5, 1.95, 6.0, 4.7, fill=LIGHT_BG)
add_bullets(s, [
    "Endorse the Hybrid identity model (MI reads + OBO writes)",
    "Endorse PIM Enablement scope as defined here",
    "Endorse V1 scope: lane A only, 2\u20133 actions, Dev subscription",
    "Sequence the milestones: PIM mid-June, Ops V1 demo end-of-June",
    "Assign owners for the 3 gaps:",
    "    \u2022  PIM trigger mechanism",
    "    \u2022  Teams output mechanism",
    "    \u2022  PIM MCP connector reality check",
    "Confirm Arturo \u2192 PIM Enablement, Jiban \u2192 Ops V1",
], 0.7, 2.1, 5.6, 4.5, size=12)

# Right: actions on Zafin
add_text(s, "Asks of Zafin (next customer session)", 6.8, 1.5, 6.0, 0.4, size=16, bold=True, color=AZURE_BLUE)
add_card(s, 6.8, 1.95, 6.0, 4.7, fill=LIGHT_BG)
add_bullets(s, [
    "Top 2\u20133 high-volume Dev ticket categories \u2014 V1 allowlist input",
    "PIM validation criteria \u2014 the rule table the agent will execute",
    "Compliance evidence requirements \u2014 audit format + retention",
    "Teams channel(s) for Adaptive Card output",
    "Confirm Object IDs (not Client IDs) used in role assignments",
], 7.0, 2.1, 5.6, 4.5, size=12)

# Bottom takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5))
tk.fill.solid(); tk.fill.fore_color.rgb = AZURE_BLUE; tk.line.fill.background()
add_text(s,
         "Goal: walk out of this session with a single, sequenced, identity-aligned plan to deliver both demos by end of June.",
         0.7, 6.93, 11.9, 0.4, size=13, bold=True, color=WHITE)
add_footer(s, 8, TOTAL)
add_notes(s,
          "Clear asks of internal team and clear asks of Zafin. The single takeaway: leave with one "
          "sequenced, identity-aligned plan \u2014 no two P1s on the same date.")

# Save
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}")
