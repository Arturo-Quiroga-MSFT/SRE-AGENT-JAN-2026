"""Build a 10-minute talk deck on the Zafin engagement.

Output: partner-context/ZAFIN_10MIN_TALK_APR-MAY2026.pptx

Run:  python scripts/build_zafin_talk_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- Theme ---------------------------------------------------------------
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
SUBTLE = RGBColor(0x60, 0x60, 0x60)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_AMBER = RGBColor(0xCA, 0x5A, 0x00)
LIGHT_BG = RGBColor(0xF3, 0xF7, 0xFB)

OUT_PATH = (
    Path(__file__).resolve().parent.parent
    / "partner-context"
    / "ZAFIN_10MIN_TALK_APR-MAY2026.pptx"
)


# ---- Helpers -------------------------------------------------------------
def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb


def add_bullets(
    slide,
    bullets: list,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 16,
    color: RGBColor = DARK_TEXT,
    line_spacing: float = 1.25,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return tb


def add_header_bar(slide, title: str, kicker: str = ""):
    # Top color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.55))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZURE_BLUE
    # Title
    add_text(slide, title, 0.5, 0.65, 12.0, 0.6, size=28, bold=True, color=DARK_TEXT)
    if kicker:
        add_text(slide, kicker, 0.5, 1.15, 12.0, 0.4, size=14, color=SUBTLE)


def add_footer(slide, idx: int, total: int):
    add_text(
        slide,
        "Arturo Quiroga  ·  PSA  ·  Zafin engagement  ·  Apr 2026",
        0.5,
        7.05,
        9.5,
        0.3,
        size=10,
        color=SUBTLE,
    )
    add_text(
        slide,
        f"{idx} / {total}",
        12.0,
        7.05,
        1.0,
        0.3,
        size=10,
        color=SUBTLE,
        align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, notes: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


# ---- Build ---------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 7

# === Slide 1: Title =======================================================
s = add_blank(prs)
# Big colored sidebar
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.2), prs.slide_height)
side.line.fill.background()
side.fill.solid()
side.fill.fore_color.rgb = AZURE_BLUE
add_text(s, "ZAFIN", 0.5, 1.4, 3.5, 0.8, size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, "Engagement", 0.5, 2.2, 3.5, 0.6, size=24, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(
    s,
    "Azure SRE Agent on hardened\nenterprise AKS",
    4.7,
    1.6,
    8.0,
    1.6,
    size=32,
    bold=True,
)
add_text(s, "10-minute team talk", 4.7, 3.4, 8.0, 0.5, size=20, color=SUBTLE)
add_text(
    s,
    "Arturo Quiroga  ·  Partner Solutions Architect  ·  April 2026",
    4.7,
    6.5,
    8.0,
    0.4,
    size=14,
    color=SUBTLE,
)
add_notes(
    s,
    "Hi everyone. Quick 10 minutes on the Zafin engagement: what they asked for, "
    "what I delivered Jan through April, the major architectural blocker I removed, "
    "what I tested, and what's next.",
)

# === Slide 2: Partner ask =================================================
s = add_blank(prs)
add_header_bar(
    s,
    "The partner ask",
    "Zafin · George Mathew (VP CloudOps), Jijo Lawrence (SRE Lead), Zoya (Eng)",
)
add_text(
    s,
    "Goal: cut incident triage from hours → minutes on Zafin's production AKS",
    0.5,
    1.7,
    12.3,
    0.5,
    size=18,
    bold=True,
    color=AZURE_BLUE,
)
add_bullets(
    s,
    [
        "Reactive: auto-RCA, ticket enrichment, mitigation suggestions",
        "Proactive: health sweeps, predictive findings, Action Tickets",
        "Troubleshooting Cockpit: parameterized Grafana per customer / namespace / time",
        "IAM/PIM agent (out of scope for this PoC)",
    ],
    0.5,
    2.3,
    12.3,
    2.5,
    size=16,
)
# Constraint callout
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6))
box.line.color.rgb = ACCENT_AMBER
box.line.width = Pt(1.5)
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
add_text(s, "Production constraint", 0.8, 5.1, 12.0, 0.4, size=14, bold=True, color=ACCENT_AMBER)
add_text(
    s,
    "Private AKS · disabled local accounts · Azure RBAC enforced · Log Analytics + Jira (no Loki)",
    0.8,
    5.55,
    12.0,
    0.5,
    size=16,
)
add_text(
    s,
    "Pre-engagement assumption: \"SRE Agent does not support AKS behind private VNet.\"",
    0.8,
    6.05,
    12.0,
    0.5,
    size=14,
    color=SUBTLE,
)
add_footer(s, 2, TOTAL)
add_notes(
    s,
    "Zafin runs a regulated banking SaaS. Their AKS clusters are private, no local "
    "accounts, Azure RBAC enforced. The original architectural belief was that SRE "
    "Agent simply could not support that topology. That framed everything.",
)

# === Slide 3: Work delivered ==============================================
s = add_blank(prs)
add_header_bar(s, "What I delivered", "January → April 2026")

rows = [
    ("Jan", "January PoC", "Grocery demo app · Loki/Grafana/Jira MCP servers · DiagnosticExpert sub-agent · 2 working scenarios"),
    ("Mar 16", "Gap analysis", "Mapped Jan PoC → Zafin's full 2026 SRE roadmap; agreed P1–P3 work split with Ankit"),
    ("Mar 30", "Private AKS testbed", "Built aks-private-testbed/ matching Zafin's topology · 4/4 validation tests PASS"),
    ("Apr",   "Validated design pattern", "Documented architecture for production rollout (ARM control plane, no VPN/Bastion)"),
    ("Apr 14", "Locked-down RBAC testbed", "Reproduced Zoya's exact Terraform + RBAC; root-caused her kubectl failures"),
    ("Apr 28", "Pre-flight role review", "Confirmed agent MI roles are read-only; flagged 2 gaps before her smoke test"),
]
top = 1.7
row_h = 0.78
for i, (when, what, detail) in enumerate(rows):
    y = top + i * row_h
    # date pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y + 0.05), Inches(1.1), Inches(0.5))
    pill.fill.solid()
    pill.fill.fore_color.rgb = AZURE_BLUE
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = when
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = "Segoe UI"
    add_text(s, what, 1.8, y, 3.5, 0.5, size=15, bold=True)
    add_text(s, detail, 5.4, y, 7.6, 0.7, size=13, color=SUBTLE)
add_footer(s, 3, TOTAL)
add_notes(
    s,
    "Six concrete deliverables, building on each other. The Jan PoC gave us a "
    "working baseline. The March gap analysis turned the customer's roadmap into "
    "a P1-P3 backlog. The two testbeds (private + locked-down) are what unlocked "
    "the real production answer.",
)

# === Slide 4: Blocker removed =============================================
s = add_blank(prs)
add_header_bar(s, "Blocker removed", "The headline outcome")

# Big quote-style callout
q = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.6))
q.fill.solid(); q.fill.fore_color.rgb = LIGHT_BG
q.line.color.rgb = AZURE_BLUE; q.line.width = Pt(1.5)
add_text(s, "Before", 0.8, 1.85, 2.0, 0.4, size=12, bold=True, color=SUBTLE)
add_text(
    s,
    "\u201CSRE Agent does not support AKS behind private VNet\u201D  → in Zafin's ADR",
    0.8,
    2.25,
    11.8,
    0.5,
    size=18,
)
add_text(s, "After (validated Mar 30)", 0.8, 2.75, 5.0, 0.4, size=12, bold=True, color=ACCENT_GREEN)
add_text(
    s,
    "Agent reaches private clusters via ARM + Log Analytics + az aks command invoke. No VPN. No Bastion.",
    0.8,
    3.0,
    11.8,
    0.5,
    size=16,
    bold=True,
    color=ACCENT_GREEN,
)

# Two-finding panel
add_text(s, "And two follow-on findings shipped to the partner:", 0.5, 3.7, 12.3, 0.4, size=14, bold=True)
add_bullets(
    s,
    [
        "App ID vs SP Object ID gotcha — `az role assignment list --assignee` doesn't auto-resolve managed-identity client IDs. Likely root cause of Zoya's reported kubectl failures.",
        "Two-layer RBAC required — both runCommand/action (ARM) AND AKS RBAC Reader (data plane). The built-in Cluster User role does NOT include runCommand/action.",
        "AKS RBAC Reader is missing pods/log/read — kubectl logs will be Forbidden; recommended ContainerLogV2 in Log Analytics instead.",
    ],
    0.5,
    4.1,
    12.3,
    2.6,
    size=14,
)
add_footer(s, 4, TOTAL)
add_notes(
    s,
    "This is the headline. The architectural ADR said it couldn't be done. We "
    "proved it can — entirely through Azure's public control plane. Then we "
    "documented the three specific RBAC traps that were biting Zoya. That "
    "conversation went from 'blocked' to 'good to go' in about three weeks.",
)

# === Slide 5: Testing I did ==============================================
s = add_blank(prs)
add_header_bar(s, "Testing I did", "Two purpose-built testbeds + one customer pre-flight")

# Three column cards
cards = [
    (
        "aks-private-testbed/",
        "Private API server, no public endpoint",
        [
            "ARM-level cluster reads",
            "Log Analytics KQL queries",
            "az aks command invoke",
            "kubectl through tunneled command pod",
        ],
        "4 / 4 PASS  ·  Mar 30",
        ACCENT_GREEN,
    ),
    (
        "aks-locked-testbed/",
        "Mirrors Zoya's exact Terraform + RBAC",
        [
            "Reader + Custom Command Invoke",
            "AKS RBAC Reader at cluster scope",
            "Read-only matrix verified",
            "Write attempts correctly denied",
        ],
        "Findings doc shipped  ·  Apr 14",
        ACCENT_GREEN,
    ),
    (
        "agent-mi-roles.md review",
        "Pre-flight check before Zoya's smoke test",
        [
            "All 4 roles confirmed read-only",
            "2 gaps flagged proactively",
            "Object-ID verification step",
            "Suggested smoke-test prompts",
        ],
        "Cleared for testing  ·  Apr 28",
        AZURE_BLUE,
    ),
]
card_w = 4.1
gap = 0.2
for i, (title, sub, bullets, status, color) in enumerate(cards):
    x = 0.5 + i * (card_w + gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.7), Inches(card_w), Inches(5.2))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    add_text(s, title, x + 0.2, 1.85, card_w - 0.4, 0.4, size=15, bold=True, color=color)
    add_text(s, sub, x + 0.2, 2.25, card_w - 0.4, 0.5, size=12, color=SUBTLE)
    add_bullets(s, bullets, x + 0.2, 2.85, card_w - 0.4, 3.0, size=12)
    # status pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(6.3), Inches(card_w - 0.4), Inches(0.45))
    pill.fill.solid(); pill.fill.fore_color.rgb = color; pill.line.fill.background()
    tf = pill.text_frame; tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = status
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = "Segoe UI"
add_footer(s, 5, TOTAL)
add_notes(
    s,
    "Three pieces of testing. The private testbed proved the architecture. The "
    "locked-down testbed reproduced the exact Zafin RBAC config so I could "
    "root-cause Zoya's failures hands-on. And the role review on Apr 28 was a "
    "pre-flight before her smoke test — caught two gaps that would have wasted a "
    "debug cycle.",
)

# === Slide 6: May 2026 — PIM Enablement validated ========================
s = add_blank(prs)
add_header_bar(s, "May 2026  \u2014  PIM Enablement validated end-to-end",
               "Hybrid Enterprise MCP + custom gap-filler `pim-mcp` working in `MngEnvMCAP094150`")

# Quick story strip
q = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.4))
q.fill.solid(); q.fill.fore_color.rgb = LIGHT_BG
q.line.color.rgb = AZURE_BLUE; q.line.width = Pt(1.5)
add_text(s, "The story", 0.8, 1.85, 4.0, 0.4, size=12, bold=True, color=SUBTLE)
add_text(
    s,
    "Microsoft shipped the Enterprise MCP Server (preview) covering ~90% of PIM Graph reads. "
    "One critical endpoint \u2014 the one exposing pending PIM activation requests \u2014 is unreachable through it. "
    "We built a single-tool gap-filler MCP server in Container Apps with Managed Identity and proved the full chain works.",
    0.8, 2.25, 11.8, 0.85, size=14,
)

# Two-column: what works / bugs found
add_text(s, "What works today", 0.5, 3.3, 6.0, 0.4, size=15, bold=True, color=ACCENT_GREEN)
g = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.75), Inches(6.0), Inches(2.9))
g.fill.solid(); g.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
g.line.color.rgb = ACCENT_GREEN; g.line.width = Pt(1.5)
add_bullets(s, [
    "`pim-mcp` Container App live (image 0.2.4)",
    "SSE \u2192 FastMCP \u2192 MI \u2192 Graph \u2192 HTTP 200",
    "`list_pending_pim_requests` returns normalized JSON",
    "Three reproducible PowerShell scripts (assign / configure / trigger)",
    "Step-by-step REPRODUCE.md for Zafin to run in their tenant",
], 0.7, 3.9, 5.6, 2.7, size=13)

add_text(s, "Upstream bugs uncovered", 6.8, 3.3, 6.0, 0.4, size=15, bold=True, color=ACCENT_AMBER)
a = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.75), Inches(6.0), Inches(2.9))
a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
a.line.color.rgb = ACCENT_AMBER; a.line.width = Pt(1.5)
add_bullets(s, [
    "BUG-001: Graph runtime requires `*.ReadWrite.Directory` for read-only LIST (docs say Read.Directory works)",
    "BUG-002: Graph rejects $filter, $orderby AND $expand on this collection",
    "BUG-003: Enterprise MCP only mirrors `MCP.*Read.*` scopes \u2014 cannot satisfy BUG-001",
    "App-role propagation lag of 5\u201360+ minutes after grant",
], 7.0, 3.9, 5.6, 2.7, size=12)

# Bottom takeaway pill
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.45))
tk.fill.solid(); tk.fill.fore_color.rgb = AZURE_BLUE; tk.line.fill.background()
add_text(
    s,
    "Same pattern as the AKS work: meet the customer where they are, document the bugs we hit, ship a defensible path forward.",
    0.7, 6.83, 11.9, 0.4, size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
)
add_footer(s, 6, TOTAL)
add_notes(
    s,
    "This slide is the May update on top of the April story. The Zafin PIM Enablement use "
    "case is now validated end-to-end in our test tenant. We picked up three real upstream "
    "bugs in the process and documented them in UPSTREAM_BUGS.md so when Zafin runs the same "
    "flow they don't re-hit the same walls. Same playbook as the AKS work \u2014 prove it, document it, "
    "ship it.",
)

# === Slide 7: What's next + close ========================================
s = add_blank(prs)
add_header_bar(s, "What's next + adjacent value", "May 2026 →")

# Left: next steps
add_text(s, "Zafin — in flight", 0.5, 1.7, 6.0, 0.4, size=16, bold=True, color=AZURE_BLUE)
add_bullets(
    s,
    [
        "P1 (me): DiagnosticExpert enriches existing Jira tickets + parameterized Grafana cockpit",
        "P1 (Ankit): Infra Debug Bundle sub-agent on KQL",
        "P2 (me): Orchestrator/Classifier sub-agent",
        "Shared: Confluence MCP for RCA Agent",
    ],
    0.5,
    2.2,
    6.0,
    3.5,
    size=13,
)

# Right: blog work
add_text(s, "Adjacent — blog with Deepthi's team", 6.8, 1.7, 6.0, 0.4, size=16, bold=True, color=AZURE_BLUE)
add_bullets(
    s,
    [
        "Draft: \"Making Azure SRE Agent Work on Locked-Down AKS Clusters\"",
        "Scope expanded per Deepthi's feedback to also cover locked-down ACA",
        "ACA testbed validated end-to-end (Apr 16) — full pass",
        "CELA review in progress; publish on Apps on Azure blog",
    ],
    6.8,
    2.2,
    6.0,
    3.5,
    size=13,
)

# Bottom takeaway
tk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0))
tk.fill.solid(); tk.fill.fore_color.rgb = AZURE_BLUE; tk.line.fill.background()
add_text(
    s,
    "Takeaway: turned a \"can't be done\" architectural blocker into a validated pattern,"
    " a confirmed customer rollout path, and a public blog so other partners benefit.",
    0.7,
    5.95,
    11.9,
    0.7,
    size=15,
    bold=True,
    color=RGBColor(0xFF, 0xFF, 0xFF),
)
add_footer(s, 7, TOTAL)
add_notes(
    s,
    "Two threads going forward. Zafin: I keep the P1 ticket-enrichment + Grafana "
    "cockpit work moving with Ankit. Adjacent: Deepthi's team asked me to "
    "productize what we learned as a public blog, expanded to ACA per her "
    "feedback. Net: one engagement, three layers of value — partner unblocked, "
    "internal pattern documented, public asset shipped.",
)

# Save
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}")
