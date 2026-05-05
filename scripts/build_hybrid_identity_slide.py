"""Build a one-slide PPTX: Hybrid identity recommendation for Zafin agents.

Output: partner-context/ZAFIN_HYBRID_IDENTITY_ONESLIDE.pptx

Run:  python scripts/build_hybrid_identity_slide.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- Theme ---------------------------------------------------------------
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
SUBTLE = RGBColor(0x60, 0x60, 0x60)
GREEN = RGBColor(0x10, 0x7C, 0x10)
AMBER = RGBColor(0xCA, 0x5A, 0x00)
RED = RGBColor(0xC4, 0x31, 0x4B)
LIGHT_BG = RGBColor(0xF3, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUT_PATH = (
    Path(__file__).resolve().parent.parent
    / "partner-context"
    / "ZAFIN_HYBRID_IDENTITY_ONESLIDE.pptx"
)


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size=14,
    bold=False,
    color=DARK_TEXT,
    align=PP_ALIGN.LEFT,
    name="Segoe UI",
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = name
    return tb


def add_bullets(slide, items, left, top, width, height, *, size=12, color=DARK_TEXT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = f"•  {it}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
    return tb


def add_card(slide, x, y, w, h, *, fill=WHITE, border=AZURE_BLUE, border_pt=1.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = border; card.line.width = Pt(border_pt)
    return card


def add_pill(slide, x, y, w, h, text, *, fill, text_color=WHITE, size=11):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    pill.fill.solid(); pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = text_color; r.font.name = "Segoe UI"
    return pill


# ---- Build ---------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Top bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.5))
bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = AZURE_BLUE

# Title
add_text(s, "Agent identity model — recommendation", 0.5, 0.6, 12.5, 0.55, size=26, bold=True)
add_text(
    s,
    "Hybrid: Managed Identity for reads, On-Behalf-Of for writes  ·  Defensible attribution for banking-regulated workloads",
    0.5, 1.15, 12.5, 0.4, size=13, color=SUBTLE,
)

# ---- 3-column comparison cards ------------------------------------------
card_w = 4.1
card_h = 3.6
gap = 0.2
y = 1.75

cards = [
    {
        "title": "Managed Identity only",
        "subtitle": "Agent acts with its own elevated rights",
        "pros": [
            "Simple to implement",
            "Works in async / scheduled flows",
            "Clean MI audit trail",
        ],
        "cons": [
            "\"Confused deputy\" risk — agent can act with no user trigger",
            "Privilege held permanently, not bound to a request",
            "Hard to attribute action to a human in audit logs",
        ],
        "verdict": "Not acceptable for banking",
        "verdict_color": RED,
        "border": SUBTLE,
    },
    {
        "title": "On-Behalf-Of only",
        "subtitle": "Agent acts under requesting user's token",
        "pros": [
            "Inherits user RBAC — no escalation",
            "Action attributed to the human",
            "Clean compliance story",
        ],
        "cons": [
            "Async / scheduled flows fail (no user token)",
            "Token lifetime issues for long-running tasks",
            "Doesn't work for proactive sweeps",
        ],
        "verdict": "Partial — writes only",
        "verdict_color": AMBER,
        "border": SUBTLE,
    },
    {
        "title": "Hybrid  (recommended)",
        "subtitle": "MI for reads, OBO for writes",
        "pros": [
            "Reads stay on validated MI pattern (Apr 14 testbed)",
            "Every write carries human attribution via OBO",
            "Defensible audit trail per action",
        ],
        "cons": [
            "Slightly more implementation complexity",
            "User identity must be present in the request chain",
        ],
        "verdict": "Recommended for banking compliance",
        "verdict_color": GREEN,
        "border": GREEN,
    },
]

for i, c in enumerate(cards):
    x = 0.5 + i * (card_w + gap)
    add_card(s, x, y, card_w, card_h, border=c["border"], border_pt=2.0 if i == 2 else 1.0)
    add_text(s, c["title"], x + 0.2, y + 0.15, card_w - 0.4, 0.4, size=15, bold=True, color=c["border"] if i == 2 else DARK_TEXT)
    add_text(s, c["subtitle"], x + 0.2, y + 0.55, card_w - 0.4, 0.4, size=11, color=SUBTLE)
    # Pros
    add_text(s, "Pros", x + 0.2, y + 1.0, card_w - 0.4, 0.3, size=11, bold=True, color=GREEN)
    add_bullets(s, c["pros"], x + 0.2, y + 1.3, card_w - 0.4, 1.0, size=10)
    # Cons
    add_text(s, "Cons", x + 0.2, y + 2.35, card_w - 0.4, 0.3, size=11, bold=True, color=RED)
    add_bullets(s, c["cons"], x + 0.2, y + 2.65, card_w - 0.4, 1.0, size=10)

# Verdict pills under cards
for i, c in enumerate(cards):
    x = 0.5 + i * (card_w + gap)
    add_pill(s, x, y + card_h + 0.1, card_w, 0.4, c["verdict"], fill=c["verdict_color"], size=12)

# ---- Bottom: per-use-case mapping ---------------------------------------
yb = 6.05
add_card(s, 0.5, yb, 12.3, 1.05, fill=LIGHT_BG, border=AZURE_BLUE, border_pt=1.0)
add_text(s, "Applied to Zafin's two priority use cases", 0.7, yb + 0.07, 12.0, 0.35, size=13, bold=True, color=AZURE_BLUE)

mapping = [
    ("PIM Enablement", "MI only (read-only end-to-end)", "Agent never writes; recommends only. Human approves in PIM portal.", GREEN),
    ("Ops Automation V1", "MI for reads + OBO for write actions", "Restart / scale / rollback — executed under requesting user's token. Dev subscription only.", GREEN),
]
col_w = 6.0
for i, (uc, model, why, color) in enumerate(mapping):
    x = 0.7 + i * (col_w + 0.1)
    add_text(s, uc, x, yb + 0.45, col_w, 0.3, size=12, bold=True, color=DARK_TEXT)
    add_text(s, model, x, yb + 0.7, col_w, 0.25, size=11, bold=True, color=color)
    add_text(s, why, x, yb + 0.92, col_w, 0.35, size=10, color=SUBTLE)

# May 5 update strip (above footer)
yu = 6.78
upd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(yu), Inches(12.3), Inches(0.32))
upd.fill.solid(); upd.fill.fore_color.rgb = LIGHT_BG
upd.line.color.rgb = AMBER; upd.line.width = Pt(1.0)
add_text(
    s,
    "Update 2026-05-05  \u00b7  PIM Enablement gap-filler MI now also holds RoleAssignmentSchedule.ReadWrite.Directory at runtime "
    "(Graph BUG-001 workaround). Latent only \u2014 server registers no write tool. See pim-enablement-testbed/docs/UPSTREAM_BUGS.md.",
    0.7, yu + 0.04, 11.9, 0.28, size=9, color=DARK_TEXT,
)

# Footer
add_text(
    s,
    "Arturo Quiroga  \u00b7  PSA  \u00b7  Zafin internal alignment  \u00b7  May 1, 2026  \u00b7  rev May 5",
    0.5, 7.18, 12.5, 0.25, size=9, color=SUBTLE,
)

# Speaker notes
notes = (
    "One-slide recommendation for the internal meeting. "
    "MI-only is off the table for banking — no per-action human attribution. "
    "OBO-only breaks proactive and scheduled flows. "
    "Hybrid is the only model that satisfies both constraints: keep the validated read pattern, "
    "and force every write to carry the requesting user's identity. "
    "Applied to our two priorities: PIM Enablement stays MI-only because it's read-only end-to-end. "
    "Ops Automation V1 uses MI for diagnostics reads and OBO for the 2–3 allowlisted write actions, "
    "scoped to the Dev subscription. This is the line I'd defend in front of a banking auditor."
)
s.notes_slide.notes_text_frame.text = notes

OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}")
